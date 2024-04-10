import tkinter as tk
from tkinter import ttk
from pptx import Presentation


num_sl = 0

def registration():   
    print('Заголовок =', e_title.get(),'Текст =',e_subtitle.get("1.0",'end-1c') )
    slide = rootp.slides[int(cbox_sl.get()) - 1 ]
    print('Номер слайда: ', int( cbox_sl.get()) - 1 )
    slide.shapes.title.text = e_title.get()
    slide.placeholders[1].text = e_subtitle.get("1.0",'end-1c')
    return
       
def new_sl():
    slide = rootp.slides.add_slide(first_slide_layout) #объект слайда
    slide.shapes.title.text = e_title.get()
    slide.placeholders[1].text = e_subtitle.get("1.0",'end-1c')
    global num_sl
    num_sl += 1
    if num_sl + 1 not in cbox_sl['values']: #добавить значение в combobox
       cbox_sl['values'] += (num_sl + 1)
    cbox_sl.current(num_sl)
    print(num_sl)
    
rootp = Presentation()
first_slide_layout = rootp.slide_layouts[1] #макет слайда
slide = rootp.slides.add_slide(first_slide_layout) #объект слайда
slide.shapes.title.text = 'Тестовый слайд' #заголовок
slide.placeholders[1].text = 'Текст тестового слайда' #текст

root = tk.Tk()
#root.geometry("500x300")
root.title('Курсовая')

lb1 = tk.Label(root, text="Заголовок", font=('Arial',18,'normal'),padx= 25, pady = 15)
e_title = tk.Entry(root,font=('Arial',18,'normal'), width= 50 )
lb2 = tk.Label(root, text="Текст", font=('Arial',18,'normal'))
e_subtitle = tk.Text(root,font=('Arial',18,'normal'), width= 50, height= 10 )
reg_button = tk.Button(root, text="Применить", font=('Arial',14,'normal'), command=registration)
add_sl = tk.Button(root, text="Новый слайд", font=('Arial',14,'normal'), command=new_sl)
cbox_sl = ttk.Combobox(root,values = 1)
cbox_sl.current(0)


lb1.grid(row=0,column=0) #размещение
e_title.grid(row=0,column=1)
lb2.grid(row=1,column=0)
e_subtitle.grid(row=1,column=1)
reg_button.grid(row=2, column=0, columnspan=2, stick='we',padx=30,pady=10)
add_sl.grid(row=3, column=0, columnspan=2, stick='we',padx=30,pady=10)
cbox_sl.grid(row=4, column=0, columnspan=2, stick='we',padx=30,pady=10)


path = r'c:/Users/user/Documents/MyCodePython/test.pptx'
root.mainloop()
rootp.save(path)